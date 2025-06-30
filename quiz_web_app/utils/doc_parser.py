import io
from typing import Optional
from fastapi import UploadFile
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation


async def parse_file(upload_file: UploadFile) -> Optional[str]:
    """Extract text from an uploaded file."""
    ext = upload_file.filename.split('.')[-1].lower()

    data = await upload_file.read()
    if ext == 'pdf':
        doc = fitz.open(stream=data, filetype='pdf')
        texts = [page.get_text() for page in doc]
        return "\n".join(texts)
    elif ext == 'docx':
        file_obj = io.BytesIO(data)
        doc = Document(file_obj)
        texts = [p.text for p in doc.paragraphs]
        return "\n".join(texts)
    elif ext == 'pptx':
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    else:
        return None
